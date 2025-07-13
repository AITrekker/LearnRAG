import asyncio
from pathlib import Path
from typing import Optional

import PyPDF2
import docx
import openpyxl
from pptx import Presentation


class FileProcessor:
    """Service for extracting text content from various file types"""

    async def extract_text(self, file_path: Path, content_type: str) -> str:
        """Extract text content from file based on content type"""
        try:
            if content_type == 'text/plain':
                return await self._extract_text_file(file_path)
            elif content_type == 'application/pdf':
                return await self._extract_pdf(file_path)
            elif content_type in [
                'application/msword',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            ]:
                return await self._extract_docx(file_path)
            elif content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            ]:
                return await self._extract_excel(file_path)
            elif content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            ]:
                return await self._extract_powerpoint(file_path)
            else:
                print(f"Unsupported content type: {content_type}")
                return ""
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""

    async def _extract_text_file(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    async def _extract_docx(self, file_path: Path) -> str:
        """Extract text from Word document"""
        if file_path.suffix.lower() == '.docx':
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        else:
            # For .doc files, we'd need python-docx2txt or similar
            # For now, return empty string
            print(f"Legacy .doc format not supported: {file_path}")
            return ""

    async def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel file"""
        if file_path.suffix.lower() in ['.xlsx', '.xlsm']:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                text += "\n"
            return text
        else:
            # For .xls files, we'd need xlrd or similar
            print(f"Legacy .xls format not supported: {file_path}")
            return ""

    async def _extract_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint file"""
        if file_path.suffix.lower() in ['.pptx']:
            presentation = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        else:
            # For .ppt files, we'd need different library
            print(f"Legacy .ppt format not supported: {file_path}")
            return ""